import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    base_dir = os.path.dirname(os.path.abspath(__file__))
    icons_dir = os.path.join(base_dir, "icons")

    # Color Palette
    BG_COLOR = RGBColor(10, 15, 29)        # #0a0f1d
    CARD_BG = RGBColor(22, 33, 56)         # #162138
    PRIMARY = RGBColor(16, 185, 129)       # #10b981 (Emerald)
    SECONDARY = RGBColor(6, 182, 212)      # #06b6d4 (Cyan)
    ACCENT = RGBColor(245, 158, 11)        # #f59e0b (Amber)
    DANGER = RGBColor(239, 68, 68)         # #ef4444 (Coral/Red)
    TEXT_MAIN = RGBColor(248, 250, 252)    # White
    TEXT_MUTED = RGBColor(148, 163, 184)   # Gray
    CARD_BORDER = RGBColor(40, 58, 92)

    # Icon Box Tints
    TINT_EMERALD = RGBColor(18, 50, 48)
    TINT_CYAN = RGBColor(16, 48, 64)
    TINT_AMBER = RGBColor(56, 44, 20)
    TINT_CORAL = RGBColor(58, 28, 32)

    blank_slide_layout = prs.slide_layouts[6]

    def set_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()
        return bg

    def add_header(slide, tag_text, title_text):
        # Tag
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.35))
        tf = tag_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = tag_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = PRIMARY

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.75))
        tf2 = title_box.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_MAIN

    def add_card(slide, left, top, width, height, title, desc, tag=None, tag_color=PRIMARY, border_color=CARD_BORDER, icon_name=None, icon_bg=None, inline_icon=False):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)

        cur_top = top + 0.18
        text_left = left + 0.22
        text_w = width - 0.44

        if icon_name:
            icon_path = os.path.join(icons_dir, f"{icon_name}.png")
            if os.path.exists(icon_path):
                if inline_icon:
                    # Place icon next to title
                    ib = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + 0.22), Inches(top + 0.18), Inches(0.48), Inches(0.48))
                    ib.fill.solid()
                    ib.fill.fore_color.rgb = icon_bg if icon_bg else TINT_EMERALD
                    ib.line.fill.background()

                    slide.shapes.add_picture(icon_path, Inches(left + 0.28), Inches(top + 0.24), Inches(0.36), Inches(0.36))
                    text_left = left + 0.8
                    text_w = width - 1.0
                else:
                    # Place icon box on top of title
                    ib = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + 0.22), Inches(top + 0.18), Inches(0.52), Inches(0.52))
                    ib.fill.solid()
                    ib.fill.fore_color.rgb = icon_bg if icon_bg else TINT_EMERALD
                    ib.line.fill.background()

                    slide.shapes.add_picture(icon_path, Inches(left + 0.28), Inches(top + 0.24), Inches(0.4), Inches(0.4))
                    cur_top = top + 0.8

        tb = slide.shapes.add_textbox(Inches(text_left), Inches(cur_top), Inches(text_w), Inches(height - (cur_top - top) - 0.15))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(14.5)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_MAIN
        p_title.space_after = Pt(3)

        if tag:
            p_tag = tf.add_paragraph()
            p_tag.text = f"[{tag}]"
            p_tag.font.size = Pt(9.5)
            p_tag.font.bold = True
            p_tag.font.color.rgb = tag_color
            p_tag.space_after = Pt(4)

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = TEXT_MUTED
        return card

    # ==========================================
    # SLIDE 1: Title & Hook
    # ==========================================
    s1 = prs.slides.add_slide(blank_slide_layout)
    set_bg(s1)

    tag_box = s1.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.4))
    p = tag_box.text_frame.paragraphs[0]
    p.text = "AI HACKATHON 2026 • LIVE PRESENTATION"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    h_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.05), Inches(11.7), Inches(1.1))
    p_h = h_box.text_frame.paragraphs[0]
    p_h.text = "NutriSense"
    p_h.font.size = Pt(46)
    p_h.font.bold = True
    p_h.font.color.rgb = PRIMARY

    sub_box = s1.shapes.add_textbox(Inches(0.8), Inches(2.15), Inches(11.7), Inches(0.6))
    p_sub = sub_box.text_frame.paragraphs[0]
    p_sub.text = "AI-Powered Precision Nutrition, Clinical Safety & Metabolic Health Suite for South Asia"
    p_sub.font.size = Pt(16.5)
    p_sub.font.color.rgb = TEXT_MAIN

    pres_box = s1.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(0.5))
    p_pres = pres_box.text_frame.paragraphs[0]
    p_pres.text = "Presented by:  Jamal Matloob   |   Affan Asim"
    p_pres.font.size = Pt(13.5)
    p_pres.font.bold = True
    p_pres.font.color.rgb = SECONDARY

    card_w = 2.75
    gap = 0.23
    top_pos = 3.6
    h_pos = 3.2

    add_card(s1, 0.8 + 0*(card_w+gap), top_pos, card_w, h_pos, "South Asian Vision", "Accurately decomposes composite gravies, curries, and cooking ghee/tarka.", "VISION AI", PRIMARY, CARD_BORDER, "camera_emerald", TINT_EMERALD)
    add_card(s1, 0.8 + 1*(card_w+gap), top_pos, card_w, h_pos, "3-Step Clinical RAG", "Independent risk evaluator and sub-50ms hypoglycemia interceptor.", "SAFETY GUARD", DANGER, CARD_BORDER, "shield_coral", TINT_CORAL)
    add_card(s1, 0.8 + 2*(card_w+gap), top_pos, card_w, h_pos, "Ramadan Fasting", "Dynamic Sehri/Iftar meal slots, dual countdown clocks & hydration pacing.", "CULTURE", SECONDARY, CARD_BORDER, "moon_cyan", TINT_CYAN)
    add_card(s1, 0.8 + 3*(card_w+gap), top_pos, card_w, h_pos, "Family & Urdu Voice", "1-tap dependent switcher with authentic Nastaleeq & Neural Urdu Voice TTS.", "INCLUSIVITY", ACCENT, CARD_BORDER, "users_amber", TINT_AMBER)

    # ==========================================
    # SLIDE 2: Problem
    # ==========================================
    s2 = prs.slides.add_slide(blank_slide_layout)
    set_bg(s2)
    add_header(s2, "The Challenge", "The South Asian Chronic Health & Nutrition Gap")

    c_w = 5.7
    c_h = 2.5
    add_card(s2, 0.8, 1.75, c_w, c_h, "Severe Metabolic Burden", "South Asians face disproportionately high rates of Type-2 Diabetes, Hypertension, and early CVD, often diagnosed a full decade earlier than Western benchmarks.", "CHRONIC CRISIS", DANGER, DANGER, "pulse_coral", TINT_CORAL, inline_icon=True)
    add_card(s2, 6.8, 1.75, c_w, c_h, "Composite Cooking Blindspot", "Western apps assume isolated foods (grilled chicken + rice). Traditional dishes (Biryani, Nihari, Haleem, Daal Tarka) contain hidden oils and sauces that Western databases cannot parse.", "DIET MISMATCH", ACCENT, ACCENT, "food_amber", TINT_AMBER, inline_icon=True)
    add_card(s2, 0.8, 4.5, c_w, c_h, "Hazardous Generic AI Advice", "Unconstrained ChatGPT bots lack medical guardrails—prescribing heavy Valsalva maneuvers to hypertensive patients or failing during acute low blood sugar crises.", "SAFETY HAZARD", DANGER, DANGER, "warning_coral", TINT_CORAL, inline_icon=True)
    add_card(s2, 6.8, 4.5, c_w, c_h, "Single-Device Family Reality", "In South Asian households, one smartphone is often shared to track health for elderly diabetic parents, growing children, and spouses. Single-user apps fail this dynamic.", "DEMOGRAPHIC GAP", SECONDARY, SECONDARY, "phone_cyan", TINT_CYAN, inline_icon=True)

    # ==========================================
    # SLIDE 3: Solution Overview
    # ==========================================
    s3 = prs.slides.add_slide(blank_slide_layout)
    set_bg(s3)
    add_header(s3, "Our Solution", "NutriSense: Precision Health Built for Reality")

    c3_w = 3.7
    c3_h = 5.1
    add_card(s3, 0.8, 1.75, c3_w, c3_h, "Cultural Multimodal Vision", "Powered by Google Gemini Vision calibrated for South Asian portion units (Katoris, Rotis, Tolas) and composite cooking fats.\n\n• Deconstructs mixed curries & gravies\n• Instant 1-tap live portion modifiers\n• Barcode scanning via OpenFoodFacts", "PLATE + BARCODE SCANNER", PRIMARY, PRIMARY, "camera_emerald", TINT_EMERALD)
    add_card(s3, 4.8, 1.75, c3_w, c3_h, "3-Step Clinical RAG Guard", "Strict clinical safety pipeline that intercepts medical risks before advice reaches the user.\n\n• Profile + 7-Day History RAG\n• Independent Clinical Risk Evaluator\n• Sub-50ms Hypoglycemia Interceptor (<70 mg/dL)\n• Emergency GPS Care Locator", "SUB-SECOND SAFETY CHECK", DANGER, DANGER, "stethoscope_coral", TINT_CORAL)
    add_card(s3, 8.8, 1.75, c3_w, c3_h, "Family & Regional Inclusion", "Designed for local demographic accessibility and offline durability.\n\n• 1-tap multi-dependent profile switcher\n• Full Jameel Noori Nastaleeq typography\n• Neural Urdu Voice Audio Coach\n• Offline SQLite sync with UUID v4 idempotency", "URDU VOICE & OFFLINE SYNC", SECONDARY, SECONDARY, "users_cyan", TINT_CYAN)

    # ==========================================
    # SLIDE 4: Multimodal Plate Scanner
    # ==========================================
    s4 = prs.slides.add_slide(blank_slide_layout)
    set_bg(s4)
    add_header(s4, "Feature Deep Dive", "Multimodal Plate Scanning & Dynamic Portions")

    add_card(s4, 0.8, 1.75, c_w, c_h, "Composite Food Decomposition", "Calibrated for South Asian culinary preparation. Deconstructs multi-ingredient meals (Haleem, Karahi, Pulao, Nihari) into protein, carb, fiber, and cooking fat metrics.", "GEMINI 2.5/3.7 VISION", PRIMARY, CARD_BORDER, "food_amber", TINT_AMBER, inline_icon=True)
    add_card(s4, 6.8, 1.75, c_w, c_h, "Live 1-Tap Portion Modifiers", "Non-blocking UI allows users to instantly tune quantities (e.g. 1/2 Katori, 1.5 Roti) with immediate macro, fiber, and calorie recalculations with zero friction.", "ZERO FRICTION UI", SECONDARY, CARD_BORDER, "sliders_cyan", TINT_CYAN, inline_icon=True)
    add_card(s4, 0.8, 4.5, c_w, c_h, "Barcode & Medical Allergen Check", "Mobile barcode scanning via OpenFoodFacts cross-references sodium, sugar, and allergens against the user's chronic health conditions (diabetic/hypertensive warnings).", "ALLERGEN & SODIUM FLAGS", ACCENT, CARD_BORDER, "barcode_amber", TINT_AMBER, inline_icon=True)
    add_card(s4, 6.8, 4.5, c_w, c_h, "Natural Language Meal Logger", "Allows conversational text or voice logging for complex custom meals with strict input timeouts and prompt injection sanitization.", "SANITIZED NLP PIPELINE", PRIMARY, CARD_BORDER, "chat_emerald", TINT_EMERALD, inline_icon=True)

    # ==========================================
    # SLIDE 5: Clinical Safety & Hypoglycemia
    # ==========================================
    s5 = prs.slides.add_slide(blank_slide_layout)
    set_bg(s5)
    add_header(s5, "Clinical Architecture", "3-Step Clinical RAG & Hypoglycemia Interceptor")

    flow_bg = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.75), Inches(11.7), Inches(1.4))
    flow_bg.fill.solid()
    flow_bg.fill.fore_color.rgb = CARD_BG
    flow_bg.line.color.rgb = PRIMARY

    tb_flow = s5.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(11.3), Inches(1.2))
    tf_f = tb_flow.text_frame
    tf_f.word_wrap = True
    p_f = tf_f.paragraphs[0]
    p_f.text = "[Step 1: Patient History RAG]  ➔  [Step 2: Contextual AI Coach]  ➔  [Step 3: Independent Risk Evaluator]"
    p_f.font.size = Pt(14)
    p_f.font.bold = True
    p_f.font.color.rgb = PRIMARY
    p_f.space_after = Pt(4)

    p_f2 = tf_f.add_paragraph()
    p_f2.text = "Retrieves 7-day logs & medical profile  ➔  Generates personalized guidance  ➔  Validates medical safety & attaches warnings."
    p_f2.font.size = Pt(11)
    p_f2.font.color.rgb = TEXT_MUTED

    add_card(s5, 0.8, 3.4, c_w, 3.4, "Sub-50ms Hypoglycemia Interceptor", "Deterministic regex & heuristic scanner immediately detects critical low blood sugar (<70 mg/dL) or acute distress, halts chat, and displays emergency fast-acting carb protocols.", "EMERGENCY PROTOCOL", DANGER, DANGER, "zap_coral", TINT_CORAL, inline_icon=True)
    add_card(s5, 6.8, 3.4, c_w, 3.4, "Integrated Emergency Care Locator", "GPS-based nearest hospital finder with automatic offline fallback to Google Maps deep-links for immediate hospital routing during medical emergencies.", "GPS CARE LOCATOR", PRIMARY, PRIMARY, "mappin_emerald", TINT_EMERALD, inline_icon=True)

    # ==========================================
    # SLIDE 6: Ramadan & Workout Engine
    # ==========================================
    s6 = prs.slides.add_slide(blank_slide_layout)
    set_bg(s6)
    add_header(s6, "Specialized Modules", "Ramadan Fasting Suite & Clinical Workouts")

    add_card(s6, 0.8, 1.75, c_w, 5.1, "Ramadan Fasting Engine", "• Dynamic Meal Remapping: Automatically reconfigures meal slots to Sehri, Iftar, Post-Iftar Dinner, and Taraweeh Snack.\n\n• Dual Live Clocks: Real-time Sehri & Iftar countdown clocks for Asia/Karachi timezone.\n\n• Split Night Hydration: Paces daily fluid intake across non-fasting night windows.\n\n• Fasting Health Rules: Pre-Sehri alarm reminders and suppressed daytime meal notifications.", "RAMADAN SUITE", SECONDARY, SECONDARY, "moon_cyan", TINT_CYAN, inline_icon=True)
    add_card(s6, 6.8, 1.75, c_w, 5.1, "Personalized Condition-Aware Workouts", "• Diabetic / High Glucose: Focuses on post-prandial glucose-blunting cardio and insulin-sensitizing resistance routines.\n\n• Hypertension / High BP: Strictly excludes heavy Valsalva breath-holding lifts; enforces steady Zone-2 aerobic training.\n\n• Joint Pain / Arthritis: Replaces high-impact movements with zero-impact isometric, water-friendly, and chair routines.\n\n• Offline Fallback: Deterministic routines if network is unavailable.", "CLINICAL EXERCISE", PRIMARY, PRIMARY, "dumbbell_emerald", TINT_EMERALD, inline_icon=True)

    # ==========================================
    # SLIDE 7: Family Profiles & Urdu Accessibility
    # ==========================================
    s7 = prs.slides.add_slide(blank_slide_layout)
    set_bg(s7)
    add_header(s7, "Accessibility & Demographics", "Family Profiles & Urdu Voice Inclusion")

    add_card(s7, 0.8, 1.75, c_w, 3.3, "1-Tap Dependent Switcher (خاندانی پروفائلز)", "Manage nutrition for elderly diabetic parents, growing children, and spouses on a single phone.\n\nIntelligent Macro Calculator automatically customizes calorie targets per dependent (e.g. low-glycemic for seniors vs. growth calories for children).", "[Me] [Ayesha] [Abu]", PRIMARY, CARD_BORDER, "users_emerald", TINT_EMERALD, inline_icon=True)
    add_card(s7, 6.8, 1.75, c_w, 3.3, "Neural Urdu Voice & Nastaleeq", "Full localization in authentic Jameel Noori Nastaleeq script.\n\nUrdu Audio Coach plays natural spoken Urdu voice advice for illiterate or senior family members who cannot read English or screen text.", "URDU VOICE TTS", SECONDARY, CARD_BORDER, "mic_cyan", TINT_CYAN, inline_icon=True)
    add_card(s7, 0.8, 5.3, 11.7, 1.6, "Offline-First SQLite Sync with UUID v4 Idempotency", "Local SQLite database ensures meal and water logging work seamlessly in low-connectivity areas with zero data duplication when reconnecting to the cloud.", "OFFLINE RELIABILITY", ACCENT, CARD_BORDER, "database_amber", TINT_AMBER, inline_icon=True)

    # ==========================================
    # SLIDE 8: Tech Stack & Resilience
    # ==========================================
    s8 = prs.slides.add_slide(blank_slide_layout)
    set_bg(s8)
    add_header(s8, "Engineering Rigor", "Technical Architecture & Multi-Model Resilience")

    add_card(s8, 0.8, 1.75, c3_w, 2.7, "FastAPI Backend", "Asynchronous Python backend with JWT LRU caching, Pydantic schemas, and structured clinical endpoints.", "FASTAPI + PYTHON 3.11", PRIMARY, CARD_BORDER, "server_emerald", TINT_EMERALD, inline_icon=True)
    add_card(s8, 4.8, 1.75, c3_w, 2.7, "Gemini Multi-Key Pool", "Auto-rotates across Gemini 2.5-Flash and 3.7-Flash with stateful 60s cooldowns—guaranteeing 0% rate-limit crashes during live traffic.", "100% UPTIME POOL", SECONDARY, SECONDARY, "cpu_cyan", TINT_CYAN, inline_icon=True)
    add_card(s8, 8.8, 1.75, c3_w, 2.7, "Supabase PostgreSQL", "Row Level Security (RLS) policies guaranteeing strict patient data isolation across family profiles and logs.", "POSTGRESQL + RLS", ACCENT, CARD_BORDER, "database_amber", TINT_AMBER, inline_icon=True)

    add_card(s8, 0.8, 4.7, c_w, 2.1, "Flutter Cross-Platform", "Android APK (backward compatible from Android 5.0 Lollipop to Android 15), iOS, and Web build pipelines.", "FLUTTER 3.x", PRIMARY, CARD_BORDER, "phone_cyan", TINT_CYAN, inline_icon=True)
    add_card(s8, 6.8, 4.7, c_w, 2.1, "Doctor PDF Stream Generator", "Zero-dependency pure Python binary PDF 1.4 engine generating weekly metabolic summaries ready for physicians.", "PHYSICIAN REPORTS", SECONDARY, CARD_BORDER, "food_amber", TINT_AMBER, inline_icon=True)

    # ==========================================
    # SLIDE 9: Live Demo Flow
    # ==========================================
    s9 = prs.slides.add_slide(blank_slide_layout)
    set_bg(s9)
    add_header(s9, "Product Walkthrough", "Live Demonstration Flow (90 Seconds)")

    d_w = 2.75
    d_top = 1.75
    d_h = 3.6
    add_card(s9, 0.8 + 0*(d_w+gap), d_top, d_w, d_h, "Step 1: Scan Plate", "Capture a photo of Biryani / Daal  ➔  instant macro breakdown with 1-tap live portion adjustment.", "MULTIMODAL AI", PRIMARY, PRIMARY, "camera_emerald", TINT_EMERALD)
    add_card(s9, 0.8 + 1*(d_w+gap), d_top, d_w, d_h, "Step 2: Switch Family", "Tap [Abu (Senior)]  ➔  dashboard instantly reconfigures for diabetic macro targets.", "FAMILY PROFILES", SECONDARY, SECONDARY, "users_cyan", TINT_CYAN)
    add_card(s9, 0.8 + 2*(d_w+gap), d_top, d_w, d_h, "Step 3: Safety Alert", "Log glucose 62 mg/dL  ➔  immediate clinical interceptor and emergency action modal.", "CLINICAL GUARD", DANGER, DANGER, "zap_coral", TINT_CORAL)
    add_card(s9, 0.8 + 3*(d_w+gap), d_top, d_w, d_h, "Step 4: Ramadan/Voice", "Show live Sehri/Iftar countdown clocks & play neural Urdu spoken voice coaching.", "URDU VOICE & FASTING", ACCENT, ACCENT, "moon_cyan", TINT_CYAN)

    banner = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.65), Inches(11.7), Inches(1.1))
    banner.fill.solid()
    banner.fill.fore_color.rgb = CARD_BG
    banner.line.color.rgb = PRIMARY
    tb_b = s9.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.3), Inches(0.8))
    p_b = tb_b.text_frame.paragraphs[0]
    p_b.text = "✓ Fully Configured & Verified on Android Physical Device (APK) & Live FastAPI Server"
    p_b.font.size = Pt(13)
    p_b.font.bold = True
    p_b.font.color.rgb = PRIMARY
    p_b.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 10: Market & Roadmap
    # ==========================================
    s10 = prs.slides.add_slide(blank_slide_layout)
    set_bg(s10)
    add_header(s10, "Vision & Roadmap", "Market Impact, Monetization & Roadmap")

    add_card(s10, 0.8, 1.75, c3_w, 3.2, "300M+ Market Need", "Targeted at South Asian diaspora worldwide facing alarming rates of metabolic disease with no culturally calibrated tool.", "HIGH NEED DEMOGRAPHIC", PRIMARY, CARD_BORDER, "globe_emerald", TINT_EMERALD)
    add_card(s10, 4.8, 1.75, c3_w, 3.2, "Business Model", "• B2C Freemium: Free core tracking + Premium AI coaching & Family Profiles.\n• B2B Clinical: Standardized physician reports for clinics and diagnostic labs.", "B2C + B2B MONETIZATION", SECONDARY, CARD_BORDER, "briefcase_cyan", TINT_CYAN)
    add_card(s10, 8.8, 1.75, c3_w, 3.2, "Future Roadmap", "• Continuous Glucose Monitor (CGM) Bluetooth sync.\n• Urdu speech-to-text plate logging.\n• Blood report lab OCR analysis.", "UPCOMING MILESTONES", ACCENT, CARD_BORDER, "rocket_amber", TINT_AMBER)

    close_box = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.25), Inches(11.7), Inches(1.5))
    close_box.fill.solid()
    close_box.fill.fore_color.rgb = CARD_BG
    close_box.line.color.rgb = PRIMARY

    tb_c = s10.shapes.add_textbox(Inches(1.0), Inches(5.4), Inches(11.3), Inches(1.2))
    tf_c = tb_c.text_frame
    p_c1 = tf_c.paragraphs[0]
    p_c1.text = '"Making metabolic health culturally accurate, clinically safe, and universally accessible."'
    p_c1.font.size = Pt(16)
    p_c1.font.bold = True
    p_c1.font.color.rgb = TEXT_MAIN
    p_c1.alignment = PP_ALIGN.CENTER
    p_c1.space_after = Pt(4)

    p_c2 = tf_c.add_paragraph()
    p_c2.text = "Thank you! — Open for Questions"
    p_c2.font.size = Pt(13)
    p_c2.font.bold = True
    p_c2.font.color.rgb = PRIMARY
    p_c2.alignment = PP_ALIGN.CENTER

    output_path = os.path.abspath(os.path.join(base_dir, "NutriSense_Presentation.pptx"))
    prs.save(output_path)
    print(f"Presentation with real graphical icons saved to: {output_path}")

if __name__ == "__main__":
    create_presentation()
